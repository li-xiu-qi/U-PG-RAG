import os
from pathlib import Path
import logging

import fitz
import pandas as pd
import requests
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from urllib.parse import urlparse

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class FileConvert:
    """封装所有文件读取操作。"""

    def __init__(self):
        self.md_suffix = '.md'
        self.text_suffix = ['.txt', '.text']
        self.excel_suffix = ['.xlsx', '.xls', '.csv']
        self.pdf_suffix = '.pdf'
        self.ppt_suffix = '.pptx'
        self.html_suffix = ['.html', '.htm', '.shtml', '.xhtml']
        self.word_suffix = ['.docx', '.doc']
        self.normal_suffix = [self.md_suffix
                              ] + self.text_suffix + self.excel_suffix + [
                                 self.pdf_suffix
                             ] + self.word_suffix + [self.ppt_suffix
                                                     ] + self.html_suffix

    def convert(self, filepath: str, delete_original: bool = False):
        """转换文件格式并根据参数决定是否删除原文件。"""
        logger.info(f"Starting conversion for file: {filepath}")
        if not os.path.isabs(filepath):
            raise ValueError("File path must be an absolute path.")

        file_type = self.get_type(filepath)
        if file_type == "md":
            content = self.convert_text(filepath)
            return content, filepath
        elif file_type == "pdf":
            content = self.convert_pdf(filepath)
        elif file_type == "excel":
            content = self.convert_excel(filepath)
        elif file_type == "word":
            content = self.convert_word(filepath)
        elif file_type == "ppt":
            content = self.convert_ppt(filepath)
        elif file_type == "html":
            content = self.convert_html(filepath)
        elif file_type == "text":
            content = self.convert_text(filepath)
        else:
            raise ValueError(f"Unsupported file type for file: {filepath}")

        if delete_original:
            os.remove(filepath)
            logger.info(f"Original file {filepath} deleted")

        new_filepath = os.path.splitext(filepath)[0] + '.md'
        with open(new_filepath, 'w', encoding='utf-8') as md_file:
            md_file.write(content)
        logger.info(f"Converted file saved as: {new_filepath}")

        return content, new_filepath

    def get_type(self, filepath: str):
        """根据URI后缀获取文件类型。"""
        filepath = filepath.lower()
        if filepath.endswith(self.pdf_suffix):
            return 'pdf'

        if filepath.endswith(self.md_suffix):
            return 'md'

        if filepath.endswith(self.ppt_suffix):
            return 'ppt'

        for suffix in self.text_suffix:
            if filepath.endswith(suffix):
                return 'text'

        for suffix in self.word_suffix:
            if filepath.endswith(suffix):
                return 'word'

        for suffix in self.excel_suffix:
            if filepath.endswith(suffix):
                return 'excel'

        for suffix in self.html_suffix:
            if filepath.endswith(suffix):
                return 'html'

        return None

    def save_image(self, image, image_name, directory):
        """保存图片到指定目录。"""
        images_dir = Path(directory) / 'images'
        images_dir.mkdir(parents=True, exist_ok=True)
        image_path = images_dir / image_name

        with open(image_path, 'wb') as img_file:
            img_file.write(image)
        logger.info(f"Image saved as: {image_path}")

    def convert_pdf(self, filepath: str):
        """读取PDF文件并转换为Markdown格式。"""
        text = ''
        try:
            with fitz.open(filepath) as pages:
                for page_num, page in enumerate(pages, start=1):
                    text += page.get_text()
                    for img_index, img in enumerate(page.get_images(full=True), start=1):
                        xref = img[0]
                        base_image = pages.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_name = f"page_{page_num}_img_{img_index}.png"
                        self.save_image(image_bytes, image_name, os.path.dirname(filepath))
                        text += f"![image](images/{image_name})\n\n"
        except Exception as e:
            logger.error(f"Failed to convert PDF file: {e}")
            raise RuntimeError(f"Failed to convert PDF file: {e}")
        return text

    def convert_excel(self, filepath: str):
        """读取Excel文件并转换为Markdown格式。"""
        try:
            if filepath.endswith('.csv'):
                table = pd.read_csv(filepath, on_bad_lines='skip', encoding='utf-8')
            else:
                table = pd.read_excel(filepath)
            if table is None:
                return ''
            return table.to_markdown(index=False)
        except pd.errors.ParserError as e:
            logger.error(f"Failed to parse Excel file: {e}")
            raise RuntimeError(f"Failed to parse Excel file: {e}")
        except Exception as e:
            logger.error(f"Failed to convert Excel file: {e}")
            raise RuntimeError(f"Failed to convert Excel file: {e}")

    def convert_word(self, filepath: str) -> str:
        """读取Word文档并转换为Markdown格式。"""
        try:
            doc = Document(filepath)
            markdown_content = ""
            for paragraph in doc.paragraphs:
                markdown_content += f"{paragraph.text}\n\n"
            for i, shape in enumerate(doc.inline_shapes, start=1):
                if shape.type == 3:
                    image = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
                    image_bytes = doc.part.related_parts[image].blob
                    image_name = f"image_{i}.png"
                    self.save_image(image_bytes, image_name, os.path.dirname(filepath))
                    markdown_content += f"![image](images/{image_name})\n\n"
            return markdown_content
        except Exception as e:
            logger.error(f"Failed to convert Word file: {e}")
            raise RuntimeError(f"Failed to convert Word file: {e}")

    def convert_ppt(self, filepath: str) -> str:
        """读取PPT文件并转换为Markdown格式。"""
        try:
            presentation = Presentation(filepath)
            markdown_content = ""
            for slide_num, slide in enumerate(presentation.slides, start=1):
                slide_title = slide.shapes.title.text if slide.shapes.title else "Slide Title"
                markdown_content += f"# {slide_title}\n\n"
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        markdown_content += f"{paragraph.text}\n\n"
                    if shape.shape_type == 13:
                        image = shape.image
                        image_bytes = image.blob
                        image_name = f"slide_{slide_num}_img.png"
                        self.save_image(image_bytes, image_name, os.path.dirname(filepath))
                        markdown_content += f"![image](images/{image_name})\n\n"
            return markdown_content
        except Exception as e:
            logger.error(f"Failed to convert PPT file: {e}")
            raise RuntimeError(f"Failed to convert PPT file: {e}")

    def convert_html(self, filepath: str) -> str:
        """读取HTML文件并转换为Markdown格式。"""
        try:
            with open(filepath, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file, 'html.parser')
            markdown_content = soup.get_text()
            for i, img in enumerate(soup.find_all('img'), start=1):
                img_url = img['src']
                if self.is_absolute_url(img_url):
                    image_bytes = requests.get(img_url).content
                    image_name = f"image_{i}.png"
                    self.save_image(image_bytes, image_name, os.path.dirname(filepath))
                    markdown_content += f"![image](images/{image_name})\n\n"
            return markdown_content
        except Exception as e:
            logger.error(f"Failed to convert HTML file: {e}")
            raise RuntimeError(f"Failed to convert HTML file: {e}")

    def convert_text(self, filepath: str) -> str:
        """读取文本文件并转换为Markdown格式。"""
        try:
            with open(filepath, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Failed to convert text file: {e}")
            raise RuntimeError(f"Failed to convert text file: {e}")

    def is_absolute_url(self, url: str) -> bool:
        """检查URL是否为绝对地址。"""
        return bool(urlparse(url).netloc)